"""
Generate the CircleSafe documentation deck (10 slides, 16:9).

Theme is the frontend's "Savanna Trust" design language, lifted from
frontend/tailwind.config.ts and frontend/app/globals.css:
emerald/mint/gold on deep forest, rounded-2xl white cards on a mint-tinted
surface, hairline #E4EFE9 borders, soft card shadows, pill buttons.

    python docs/make_deck.py                 # Segoe UI (installed everywhere)
    python docs/make_deck.py --exact-fonts   # Space Grotesk + Inter, as the app ships

Content is drawn from README.md and PROGRESS.md.
"""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette --
# frontend/tailwind.config.ts → theme.extend.colors
PRIMARY = "0FA968"
PRIMARY_LIGHT = "34E0A1"
PRIMARY_DARK = "0B8452"
ACCENT = "F5B301"
ACCENT_DARK = "D99A00"
INK = "07231B"
SURFACE = "F4FAF7"
LINE = "E4EFE9"
BODY = "1C2B27"
MUTED = "5C7268"
WHITE = "FFFFFF"

EXACT = "--exact-fonts" in sys.argv
HEAD = "Space Grotesk" if EXACT else "Segoe UI"          # font-heading
SANS = "Inter" if EXACT else "Segoe UI"                  # font-sans
MONO = "Consolas"

# ------------------------------------------------------------------ geometry
W, H = Inches(13.333), Inches(7.5)
M = Inches(0.75)                  # container-x
CW = W - 2 * M                    # content width

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------------- helpers
def rgb(h):
    return RGBColor.from_string(h)


def _spPr(shape):
    return shape._element.spPr


def shadow(shape, blur=34, dist=12, color="14281E", alpha=26):
    """Approximate Tailwind's shadow-card / shadow-glow."""
    sp = _spPr(shape)
    for e in sp.findall(qn("a:effectLst")):
        sp.remove(e)
    sp.append(parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{int(blur * 12700)}" dist="{int(dist * 12700)}" dir="5400000" '
        f'rotWithShape="0"><a:srgbClr val="{color}"><a:alpha val="{int(alpha * 1000)}"/>'
        "</a:srgbClr></a:outerShdw></a:effectLst>"
    ))


def no_shadow(shape):
    shape.shadow.inherit = False


def fill_alpha(shape, color, alpha):
    """Solid fill with transparency (alpha 0-100)."""
    sp = _spPr(shape)
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for e in sp.findall(qn(tag)):
            sp.remove(e)
    ln = sp.find(qn("a:ln"))
    frag = parse_xml(
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{color}"><a:alpha val="{int(alpha * 1000)}"/></a:srgbClr></a:solidFill>'
    )
    if ln is not None:
        ln.addprevious(frag)
    else:
        sp.append(frag)


def rect(slide, x, y, w, h, fill=None, radius=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    no_shadow(s)
    s.text_frame.word_wrap = True
    return s


def card(slide, x, y, w, h, fill=WHITE, radius=0.09, line=LINE, soft=True):
    """.card — rounded-2xl border border-line bg-white shadow-card"""
    s = rect(slide, x, y, w, h, fill=fill, radius=radius, line=line, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if soft:
        shadow(s)
    return s


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paras: list of dicts {t, size, bold, color, font, space_after, line, align, spc, italic}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        para.line_spacing = p.get("line", 1.25)
        if p.get("space_after") is not None:
            para.space_after = Pt(p["space_after"])
        run = para.add_run()
        run.text = p["t"]
        f = run.font
        f.name = p.get("font", SANS)
        f.size = Pt(p.get("size", 12))
        f.bold = p.get("bold", False)
        f.italic = p.get("italic", False)
        f.color.rgb = rgb(p.get("color", BODY))
        if p.get("spc"):
            f._rPr.set("spc", str(int(p["spc"] * 100)))
    return tb


def pill(slide, x, y, w, h, label, fill=PRIMARY, color=WHITE, size=10.5, bold=True,
         line=None, alpha=None, spc=0, font=None):
    """.btn — rounded-full pill."""
    s = rect(slide, x, y, w, h, fill=fill if alpha is None else None,
             radius=0.5, line=line, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if alpha is not None:
        fill_alpha(s, fill, alpha)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = font or HEAD
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    if spc:
        r.font._rPr.set("spc", str(int(spc * 100)))
    return s


def badge(slide, x, y, d, label, fill=PRIMARY, color=WHITE, size=13, alpha=None, square=True):
    """Icon-in-tinted-square, as on the frontend feature cards."""
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if square else MSO_SHAPE.OVAL
    s = rect(slide, x, y, d, d, fill=fill if alpha is None else None,
             radius=0.22 if square else None, shape=shp)
    if alpha is not None:
        fill_alpha(s, fill, alpha)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = HEAD
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    return s


def slide_bg(slide, color):
    rect(slide, 0, 0, W, H, fill=color)


def glow(slide, cx, cy, d, color, alpha):
    """Soft brand halo, echoing the Lightfall hero."""
    s = rect(slide, Emu(int(cx - d / 2)), Emu(int(cy - d / 2)), Emu(int(d)), Emu(int(d)),
             shape=MSO_SHAPE.OVAL)
    fill_alpha(s, color, alpha)
    return s


def header(slide, tag, title, sub=None, dark=False):
    """.section-tag + heading, matching the landing sections."""
    tcol = WHITE if dark else BODY
    text(slide, M, Inches(0.52), CW, Inches(0.3),
         [{"t": tag, "size": 11, "bold": True, "color": PRIMARY_LIGHT if dark else PRIMARY,
           "font": HEAD, "spc": 1.4}])
    text(slide, M, Inches(0.86), CW, Inches(0.6),
         [{"t": title, "size": 30, "bold": True, "color": tcol, "font": HEAD, "line": 1.05}])
    if sub:
        text(slide, M, Inches(1.44), Inches(9.6), Inches(0.4),
             [{"t": sub, "size": 12.5, "color": "9FB8AD" if dark else MUTED, "font": SANS,
               "line": 1.35}])


def footer(slide, n, dark=False):
    col = "6C8A7D" if dark else MUTED
    text(slide, M, H - Inches(0.62), Inches(8), Inches(0.3),
         [{"t": "CircleSafe  ·  Trustworthy Ajo. Every naira accounted for.", "size": 9, "color": col}])
    text(slide, W - M - Inches(2), H - Inches(0.62), Inches(2), Inches(0.3),
         [{"t": f"{n:02d} / 10", "size": 9, "color": col, "font": HEAD, "bold": True,
           "align": PP_ALIGN.RIGHT}])


def new_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    slide_bg(s, INK if dark else SURFACE)
    return s


def rule(slide, x, y, w, color=LINE, h=Pt(1)):
    rect(slide, x, y, w, Emu(int(h)), fill=color)


# =========================================================== 1 · title ======
s = new_slide(dark=True)
glow(s, Inches(11.2), Inches(1.0), Inches(7.2), PRIMARY, 16)
glow(s, Inches(12.6), Inches(6.8), Inches(5.2), PRIMARY_LIGHT, 12)
glow(s, Inches(0.4), Inches(7.4), Inches(4.4), ACCENT, 8)

pill(s, M, Inches(1.15), Inches(4.62), Inches(0.4),
     "ZERO TO QUERY  ·  SUB0 + LINGOQL", fill=PRIMARY_LIGHT, color=PRIMARY_LIGHT,
     alpha=14, size=10, spc=1.2)

text(s, M, Inches(1.82), Inches(9.4), Inches(1.5),
     [{"t": "CircleSafe", "size": 72, "bold": True, "color": WHITE, "font": HEAD, "line": 1.0}])
text(s, M, Inches(3.08), Inches(9.4), Inches(0.6),
     [{"t": "Trustworthy Ajo. Every naira accounted for.", "size": 21, "bold": True,
       "color": PRIMARY_LIGHT, "font": HEAD}])
text(s, M, Inches(3.75), Inches(7.9), Inches(1.1),
     [{"t": "A transparent, append-only system of record for rotating savings groups — "
            "Ajo · Esusu · Susu · Tontine — the informal circles that move enormous amounts "
            "of money across West Africa every week on nothing but a notebook and trust.",
       "size": 13, "color": "AFC7BC", "line": 1.5}])

rule(s, M, Inches(5.12), Inches(11.83), color="14392C", h=Pt(1))

STATS = [("39", "API endpoints"), ("11", "Postgres tables"), ("1", "Append-only audit log"),
         ("100%", "Auditable history")]
sx, sw = M, Inches(2.55)
for i, (v, l) in enumerate(STATS):
    x = sx + i * Inches(2.95)
    text(s, x, Inches(5.42), sw, Inches(0.55),
         [{"t": v, "size": 34, "bold": True, "color": WHITE, "font": HEAD, "line": 1.0}])
    text(s, x, Inches(5.98), sw, Inches(0.3),
         [{"t": l, "size": 10.5, "color": "8FAEA1"}])

text(s, M, Inches(6.72), Inches(9), Inches(0.3),
     [{"t": "Project documentation  ·  July 2026  ·  github.com/akaheem/Circle-Safe",
       "size": 10, "color": "6C8A7D"}])

# ========================================================== 2 · problem =====
s = new_slide()
header(s, "THE PROBLEM", "The tradition works. The bookkeeping doesn't.",
       "Ten traders each contribute ₦5,000 a week; each week the ₦50,000 pot goes to one member "
       "until everyone has taken a turn. It is one of West Africa's most widely used savings "
       "instruments — and it runs on paper.")

PROBLEMS = [
    ("01", "Fraud", "The collector disappears with the pot. There is no independent record to point to."),
    ("02", "Lost records", "The notebook is lost or the WhatsApp group is deleted; nobody can prove who paid."),
    ("03", "Disputes", "“I paid you last Tuesday.” “You didn't.” “Whose turn is it now?”"),
    ("04", "No financial history", "Years of reliable payments produce zero credit signal a lender will accept."),
]
cw, gap = Inches(2.83), Inches(0.17)
for i, (n, t, d) in enumerate(PROBLEMS):
    x = M + i * (cw + gap)
    card(s, x, Inches(2.32), cw, Inches(2.62))
    badge(s, x + Inches(0.36), Inches(2.68), Inches(0.56), n, fill=PRIMARY, alpha=12, color=PRIMARY, size=15)
    text(s, x + Inches(0.36), Inches(3.44), cw - Inches(0.72), Inches(0.4),
         [{"t": t, "size": 14.5, "bold": True, "color": BODY, "font": HEAD}])
    text(s, x + Inches(0.36), Inches(3.86), cw - Inches(0.72), Inches(1.0),
         [{"t": d, "size": 10.5, "color": MUTED, "line": 1.45}])

b = card(s, M, Inches(5.28), CW, Inches(1.32), fill=INK, line=INK)
text(s, M + Inches(0.55), Inches(5.62), CW - Inches(1.1), Inches(0.7),
     [{"t": "Four failure modes people live with today — all of them bookkeeping problems, "
            "none of them problems with the tradition itself.",
       "size": 14, "bold": True, "color": WHITE, "font": HEAD, "line": 1.35}])
footer(s, 2)

# ========================================================= 3 · solution =====
s = new_slide()
header(s, "THE SOLUTION", "Keep the tradition. Replace the notebook.",
       "Six guarantees, each enforced by the data model rather than by policy.")

SOLUTIONS = [
    ("Append-only ledger", "Every action writes a row to `_activity` that is never updated or deleted. The audit trail *is* the data model."),
    ("Two-party confirmation", "A member records a contribution; a treasurer or the owner confirms it. Neither side moves money in the record alone."),
    ("Rules locked before the money", "Amount, rhythm, grace period, late fee and payout order freeze the moment a circle goes ACTIVE — enforced in SQL."),
    ("Trust Score", "Confirmed contributions ÷ cycles elapsed, computed live from the ledger. The portable credit history the tradition never produced."),
    ("Circle Health", "One number for whether a circle is actually collecting what it is owed, so trouble is visible before it compounds."),
    ("Live dashboard", "Every mutation broadcasts — SSE self-hosted, WebSocket on Sub0 — so all members see the same numbers at the same time."),
]
cw, gap = Inches(3.83), Inches(0.19)
for i, (t, d) in enumerate(SOLUTIONS):
    col, row = i % 3, i // 3
    x = M + col * (cw + gap)
    y = Inches(2.28) + row * Inches(2.18)
    card(s, x, y, cw, Inches(1.98))
    badge(s, x + Inches(0.34), y + Inches(0.34), Inches(0.5), "✓", fill=PRIMARY, alpha=12,
          color=PRIMARY, size=14)
    text(s, x + Inches(1.0), y + Inches(0.4), cw - Inches(1.34), Inches(0.4),
         [{"t": t, "size": 14, "bold": True, "color": BODY, "font": HEAD}])
    text(s, x + Inches(0.34), y + Inches(1.0), cw - Inches(0.68), Inches(0.85),
         [{"t": d.replace("`", "").replace("*", ""), "size": 10.5, "color": MUTED, "line": 1.45}])
footer(s, 3)

# ===================================================== 4 · architecture =====
s = new_slide(dark=True)
glow(s, Inches(12.4), Inches(0.6), Inches(6.0), PRIMARY, 12)
header(s, "ARCHITECTURE", "One contract, three runtimes.",
       "The browser POSTs to a resource; the runtime validates the payload, verifies the JWT, "
       "rate-limits by IP, runs chained SQL actionables, writes the audit row and broadcasts.",
       dark=True)

# client
cy = Inches(2.42)
c1 = card(s, M, cy, Inches(3.5), Inches(2.5), fill="0C3025", line="17513C", soft=False)
text(s, M + Inches(0.34), cy + Inches(0.3), Inches(2.9), Inches(0.3),
     [{"t": "BROWSER", "size": 9.5, "bold": True, "color": PRIMARY_LIGHT, "font": HEAD, "spc": 1.2}])
text(s, M + Inches(0.34), cy + Inches(0.66), Inches(2.9), Inches(0.4),
     [{"t": "Next.js 16 · App Router", "size": 14, "bold": True, "color": WHITE, "font": HEAD}])
text(s, M + Inches(0.34), cy + Inches(1.12), Inches(2.9), Inches(1.2),
     [{"t": "React 19 · TypeScript · Tailwind\nRecharts · dnd-kit · framer-motion\nogl WebGL hero · lucide-react",
       "size": 10.5, "color": "9FB8AD", "line": 1.5}])

# arrows + labels
ax = M + Inches(3.62)
for i, (lbl, y) in enumerate([("POST /resource  ·  Bearer JWT", cy + Inches(0.72)),
                              ("wss /ws  ·  circle_update", cy + Inches(1.52))]):
    a = rect(s, ax, y, Inches(1.42), Inches(0.2), fill=PRIMARY_LIGHT if i == 0 else ACCENT,
             shape=MSO_SHAPE.RIGHT_ARROW)
    text(s, ax - Inches(0.1), y - Inches(0.3), Inches(1.7), Inches(0.25),
         [{"t": lbl, "size": 8, "color": "9FB8AD", "align": PP_ALIGN.CENTER, "font": MONO}])

# platform
px = M + Inches(5.3)
p = card(s, px, cy, Inches(6.53), Inches(2.5), fill="0C3025", line="17513C", soft=False)
text(s, px + Inches(0.34), cy + Inches(0.3), Inches(5.8), Inches(0.3),
     [{"t": "LINGOQL  ·  MANAGED PLATFORM + TLS", "size": 9.5, "bold": True,
       "color": PRIMARY_LIGHT, "font": HEAD, "spc": 1.2}])
sub = card(s, px + Inches(0.34), cy + Inches(0.7), Inches(3.0), Inches(1.5), fill="103C2E", line="1E6349", soft=False)
text(s, px + Inches(0.56), cy + Inches(0.92), Inches(2.6), Inches(1.1),
     [{"t": "Sub0", "size": 14, "bold": True, "color": WHITE, "font": HEAD, "space_after": 4},
      {"t": "Declarative JSON models\n+ ABI endpoints.\nNo hand-written server code.",
       "size": 10, "color": "9FB8AD", "line": 1.4}])
db = card(s, px + Inches(3.52), cy + Inches(0.7), Inches(2.67), Inches(1.5), fill="103C2E", line="1E6349", soft=False)
text(s, px + Inches(3.74), cy + Inches(0.92), Inches(2.3), Inches(1.1),
     [{"t": "PostgreSQL", "size": 14, "bold": True, "color": WHITE, "font": HEAD, "space_after": 4},
      {"t": "11 tables.\nEvery query parameterized.\n_activity is INSERT-only.",
       "size": 10, "color": "9FB8AD", "line": 1.4}])
a = rect(s, px + Inches(3.36), cy + Inches(1.35), Inches(0.16), Inches(0.2),
         fill=PRIMARY_LIGHT, shape=MSO_SHAPE.RIGHT_ARROW)

# runtimes strip
RUNTIMES = [
    ("mock", "Default", "In-memory fake of every resource. Clickable with no backend. Header badge reads Demo."),
    ("local", "Self-hosted", "Next.js route handlers + pg + bcryptjs + jsonwebtoken over any PostgreSQL."),
    ("sub0", "Canonical", "The declarative design in backend/ — 6 models, 23 ABI endpoints."),
]
cw = Inches(3.83)
for i, (k, tag, d) in enumerate(RUNTIMES):
    x = M + i * (cw + Inches(0.19))
    y = Inches(5.22)
    c = card(s, x, y, cw, Inches(1.28), fill="0C3025", line="17513C", soft=False)
    text(s, x + Inches(0.3), y + Inches(0.22), cw - Inches(0.6), Inches(0.3),
         [{"t": f"{k}  —  {tag}", "size": 11.5, "bold": True, "color": ACCENT, "font": MONO}])
    text(s, x + Inches(0.3), y + Inches(0.56), cw - Inches(0.6), Inches(0.6),
         [{"t": d, "size": 9.5, "color": "9FB8AD", "line": 1.4}])
text(s, M, Inches(4.98), CW, Inches(0.25),
     [{"t": "One client seam — frontend/lib/api.ts — switches all three on two environment variables.",
       "size": 10, "color": "6C8A7D", "italic": True}])
footer(s, 4, dark=True)

# ======================================================= 5 · data model =====
s = new_slide()
header(s, "DATA MODEL", "Six core tables. Nothing to falsify.",
       "Trust Score, Circle Health and Insights are computed, not stored — cached aggregation "
       "queries over _contributions and _memberships. Nothing to drift.")

ROWS = [
    ("_users", "Accounts", "email unique + indexable; password BCRYPT-hashed, never returned; Google Sign-In supported"),
    ("_circles", "The savings group", "status PENDING → ACTIVE → COMPLETED; current_cycle; rules as a JSON object"),
    ("_memberships", "Who is in a circle", "role OWNER / TREASURER / MEMBER; payout_position; status"),
    ("_contributions", "Money in", "cycle; amount as NUMERIC (exact decimal); status PENDING / CONFIRMED"),
    ("_payouts", "Money out", "cycle; recipient_id; status SCHEDULED / PAID / RECEIVED"),
    ("_activity", "Append-only audit log", "INSERT only — never UPDATE, never DELETE, by construction"),
]
tx, ty, tw = M, Inches(2.26), CW
rowh = Inches(0.56)
# head
hd = rect(s, tx, ty, tw, Inches(0.42), fill=INK, radius=0.28, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
for lbl, ox, ow in [("TABLE", Inches(0.34), Inches(2.4)), ("PURPOSE", Inches(2.9), Inches(2.4)),
                    ("NOTES", Inches(5.5), Inches(6.0))]:
    text(s, tx + ox, ty + Inches(0.11), ow, Inches(0.25),
         [{"t": lbl, "size": 9, "bold": True, "color": PRIMARY_LIGHT, "font": HEAD, "spc": 1.2}])
for i, (t, p, n) in enumerate(ROWS):
    y = ty + Inches(0.52) + i * rowh
    if i % 2 == 0:
        rect(s, tx, y, tw, rowh - Inches(0.06), fill=WHITE, radius=0.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, tx + Inches(0.34), y + Inches(0.12), Inches(2.5), Inches(0.3),
         [{"t": t, "size": 11.5, "bold": True, "color": PRIMARY_DARK, "font": MONO}])
    text(s, tx + Inches(2.9), y + Inches(0.13), Inches(2.5), Inches(0.3),
         [{"t": p, "size": 11, "bold": True, "color": BODY, "font": HEAD}])
    text(s, tx + Inches(5.5), y + Inches(0.14), Inches(6.1), Inches(0.3),
         [{"t": n, "size": 10, "color": MUTED}])

note = card(s, M, Inches(6.06), CW, Inches(0.72), fill=WHITE)
text(s, M + Inches(0.4), Inches(6.24), CW - Inches(0.8), Inches(0.4),
     [{"t": "The running app adds five more tables — _invitations, _join_requests, _email_tokens, "
            "_emails, _revoked_tokens — for the invitation lifecycle, join requests and JWT revocation.",
       "size": 10.5, "color": MUTED, "line": 1.35}])
footer(s, 5)

# ============================================================ 6 · sub0 ======
s = new_slide(dark=True)
glow(s, Inches(0.2), Inches(7.6), Inches(5.4), PRIMARY, 10)
header(s, "TECHNICAL IMPLEMENTATION", "How CircleSafe uses Sub0.",
       "Sub0 is declarative: you define models and endpoints as JSON and the platform is the "
       "backend. Everything below is a platform primitive, not application code we wrote.",
       dark=True)

PRIMS = [
    ("hashables / verify_hashables", "BCRYPT on sign-up, verified on sign-in. Passwords never appear in any returnables."),
    ("tokenize / invalidate_tokenize", "JWT HS256 signed from $ENV.JWT_SECRET_KEY; logout revokes the presented token."),
    ("protected", "Every non-public endpoint. Ownership always comes from $PROTECTED.id, never from the payload."),
    ("payload_validation", "Types and length bounds on every endpoint that accepts input — EMAIL, STRING, NUMBER."),
    ("rate_limit", "Auth, create-circle, invite-member and record-contribution, keyed on $HEADER.ip."),
    ("depends_on (action chaining)", "create-circle → owner membership → audit row; record-payout computes the pot, pays, advances the cycle."),
    ("read_from_cache", "60s on trust score, circle health and insights. get-dashboard is deliberately uncached."),
    ("broadcast_websocket_message", "Fires on all six lifecycle mutations, from start-circle to confirm-payout-received."),
    ("$GENERATOR.KSUID", "Every primary key — sortable, collision-resistant, no sequence to leak."),
    ("Models & injection prevention", "JSON_OBJECT, foreign_key, indexable; every query parameterized before execution."),
]
cw = Inches(5.82)
for i, (k, d) in enumerate(PRIMS):
    col, row = i % 2, i // 2
    x = M + col * (cw + Inches(0.19))
    y = Inches(2.3) + row * Inches(0.79)
    text(s, x, y, cw, Inches(0.28),
         [{"t": k, "size": 11, "bold": True, "color": PRIMARY_LIGHT, "font": MONO}])
    text(s, x, y + Inches(0.27), cw, Inches(0.44),
         [{"t": d, "size": 9.5, "color": "9FB8AD", "line": 1.35}])

rule(s, M, Inches(6.32), CW, color="14392C")
text(s, M, Inches(6.55), CW, Inches(0.3),
     [{"t": "23 endpoints across Auth · Circles · Members · Contributions · Payouts · Dashboard · "
            "Scores · Activity  —  the running app extends the same contract to 39 resources.",
       "size": 10, "color": "6C8A7D"}])
footer(s, 6, dark=True)

# ======================================================= 7 · lifecycle ======
s = new_slide()
header(s, "THE CORE FLOW", "No single person can move money in the record.",
       "The contribution lifecycle, end to end. Each step writes to the audit log and broadcasts "
       "to every member of the circle.")

STEPS = [
    ("01", "MEMBER", "record-contribution", "Writes to _contributions as PENDING, appends to _activity, broadcasts."),
    ("02", "TREASURER / OWNER", "confirm-contribution", "Flips the row to CONFIRMED. RBAC checked in SQL, not in the UI."),
    ("03", "OWNER", "record-payout", "Sums CONFIRMED for the cycle, pays position N, advances current_cycle."),
    ("04", "RECIPIENT", "confirm-payout-received", "Closes the loop as RECEIVED. Circle COMPLETED after the last member."),
]
cw, gap = Inches(2.83), Inches(0.17)
for i, (n, role, action, d) in enumerate(STEPS):
    x = M + i * (cw + gap)
    y = Inches(2.36)
    card(s, x, y, cw, Inches(2.5))
    badge(s, x + Inches(0.34), y + Inches(0.34), Inches(0.5), n, fill=PRIMARY, color=WHITE, size=13)
    text(s, x + Inches(0.96), y + Inches(0.47), cw - Inches(1.3), Inches(0.25),
         [{"t": role, "size": 8.5, "bold": True, "color": MUTED, "font": HEAD, "spc": 1.0}])
    text(s, x + Inches(0.34), y + Inches(1.02), cw - Inches(0.68), Inches(0.3),
         [{"t": action, "size": 11.5, "bold": True, "color": PRIMARY_DARK, "font": MONO}])
    text(s, x + Inches(0.34), y + Inches(1.42), cw - Inches(0.68), Inches(0.9),
         [{"t": d, "size": 10, "color": MUTED, "line": 1.45}])
    if i < 3:
        rect(s, x + cw + Inches(0.015), y + Inches(1.16), Inches(0.14), Inches(0.18),
             fill=PRIMARY_LIGHT, shape=MSO_SHAPE.RIGHT_ARROW)

band = card(s, M, Inches(5.22), CW, Inches(1.32), fill=WHITE)
text(s, M + Inches(0.5), Inches(5.44), Inches(3.1), Inches(0.3),
     [{"t": "CIRCLE STATUS", "size": 9, "bold": True, "color": MUTED, "font": HEAD, "spc": 1.2}])
for i, (lbl, col) in enumerate([("PENDING", ACCENT), ("ACTIVE", PRIMARY), ("COMPLETED", PRIMARY_DARK)]):
    x = M + Inches(0.5) + i * Inches(1.72)
    pill(s, x, Inches(5.8), Inches(1.4), Inches(0.36), lbl, fill=col, alpha=16, color=col, size=9.5, spc=0.8)
    if i < 2:
        rect(s, x + Inches(1.46), Inches(5.9), Inches(0.16), Inches(0.16),
             fill=LINE, shape=MSO_SHAPE.RIGHT_ARROW)
text(s, M + Inches(6.0), Inches(5.5), Inches(5.3), Inches(0.9),
     [{"t": "Rules and payout order are writable only while a circle is PENDING — enforced in the "
            "WHERE clause of the update itself. Once money moves, the terms are frozen.",
       "size": 10.5, "color": MUTED, "line": 1.45}])
footer(s, 7)

# ====================================================== 8 · innovation ======
s = new_slide()
header(s, "INNOVATION & UTILITY", "What a notebook could never do.",
       "Four things that only exist because the ledger is complete, ordered and honest.")

BIG = [
    ("Trust Score", "Confirmed contributions ÷ cycles elapsed, per member, computed live.",
     "Years of reliable payments become a portable financial history a lender can read — the single "
     "biggest gap in informal savings today."),
    ("Circle Health", "One number for whether a circle is collecting what it is owed.",
     "Surfaces a circle drifting toward collapse while there is still time to act, instead of at the "
     "cycle where the pot comes up short."),
    ("Visual Rules Builder", "Drag members into payout order; set amount, rhythm, grace and late fee.",
     "The most argued-about part of an Ajo becomes an explicit, agreed artifact — set once, frozen "
     "on start, visible to everyone."),
    ("Live shared truth", "Every mutation pushes to every member — SSE self-hosted, WebSocket on Sub0.",
     "Disputes need a single source of truth. All members see the same pot, the same cycle and the "
     "same next recipient at the same moment."),
]
cw, gap = Inches(5.82), Inches(0.19)
for i, (t, lead, d) in enumerate(BIG):
    col, row = i % 2, i // 2
    x = M + col * (cw + gap)
    y = Inches(2.28) + row * Inches(2.16)
    card(s, x, y, cw, Inches(1.96))
    accentbar = rect(s, x, y + Inches(0.34), Inches(0.06), Inches(1.28), fill=PRIMARY if i % 2 == 0 else ACCENT)
    text(s, x + Inches(0.42), y + Inches(0.32), cw - Inches(0.84), Inches(0.32),
         [{"t": t, "size": 16, "bold": True, "color": BODY, "font": HEAD}])
    text(s, x + Inches(0.42), y + Inches(0.72), cw - Inches(0.84), Inches(0.3),
         [{"t": lead, "size": 10.5, "bold": True, "color": PRIMARY_DARK, "line": 1.35}])
    text(s, x + Inches(0.42), y + Inches(1.12), cw - Inches(0.84), Inches(0.7),
         [{"t": d, "size": 10, "color": MUTED, "line": 1.45}])
footer(s, 8)

# ========================================================= 9 · the app ======
s = new_slide()
header(s, "THE APPLICATION", "Thirteen screens, one design language.",
       "Next.js App Router, grouped by access. Every screen is reachable in demo mode with no "
       "backend running.")

GROUPS = [
    ("PUBLIC", PRIMARY, [
        ("/", "Landing — hero, features, stats, FAQ"),
        ("/login  ·  /register", "Email + password, or Google Sign-In"),
        ("/invite/[token]", "Accept an emailed invitation"),
    ]),
    ("MEMBER", PRIMARY_DARK, [
        ("/dashboard", "Every circle, pot totals, next payout, live"),
        ("/circles/new", "Create a circle; the visual Rules Builder"),
        ("/circles/[id]", "Ledger, members, contributions, payouts, activity"),
        ("/circles/discover  ·  /requests", "Find public circles; manage join requests"),
    ]),
    ("ADMIN", ACCENT_DARK, [
        ("/admin", "Platform overview"),
        ("/admin/users  ·  /admin/circles", "Accounts and circles across the platform"),
        ("/admin/emails", "Outbound mail log — no provider needed to test"),
    ]),
]
cw, gap = Inches(3.83), Inches(0.19)
for i, (name, col, routes) in enumerate(GROUPS):
    x = M + i * (cw + gap)
    y = Inches(2.3)
    card(s, x, y, cw, Inches(3.3))
    pill(s, x + Inches(0.34), y + Inches(0.3), Inches(1.16), Inches(0.32), name,
         fill=col, alpha=14, color=col, size=9, spc=1.0)
    for j, (r, d) in enumerate(routes):
        ry = y + Inches(0.82) + j * Inches(0.6)
        text(s, x + Inches(0.34), ry, cw - Inches(0.68), Inches(0.25),
             [{"t": r, "size": 10, "bold": True, "color": BODY, "font": MONO}])
        text(s, x + Inches(0.34), ry + Inches(0.22), cw - Inches(0.68), Inches(0.3),
             [{"t": d, "size": 9, "color": MUTED, "line": 1.3}])

b = card(s, M, Inches(5.78), CW, Inches(0.96), fill=INK, line=INK)
text(s, M + Inches(0.5), Inches(6.0), Inches(6.4), Inches(0.55),
     [{"t": "Run it: cd frontend && npm install && npm run dev", "size": 12, "bold": True,
       "color": WHITE, "font": MONO, "space_after": 3},
      {"t": "With nothing configured the app runs in demo mode against an in-memory mock of every resource.",
       "size": 9.5, "color": "9FB8AD"}])
pill(s, W - M - Inches(4.6), Inches(6.06), Inches(2.0), Inches(0.4), "Demo  ·  mock runtime",
     fill=ACCENT, alpha=18, color=ACCENT, size=9.5)
pill(s, W - M - Inches(2.4), Inches(6.06), Inches(1.9), Inches(0.4), "Live  ·  local / sub0",
     fill=PRIMARY_LIGHT, alpha=18, color=PRIMARY_LIGHT, size=9.5)
footer(s, 9)

# ======================================================= 10 · closing =======
s = new_slide(dark=True)
glow(s, Inches(12.8), Inches(0.4), Inches(6.6), PRIMARY, 14)
glow(s, Inches(0.6), Inches(7.6), Inches(5.0), ACCENT, 7)
header(s, "STACK, SECURITY & WHAT'S NEXT", "Built to be audited.", dark=True)

COLS = [
    ("STACK", [
        "Next.js 16 · React 19 · TypeScript",
        "Tailwind CSS · Recharts · dnd-kit",
        "framer-motion · ogl · lucide-react",
        "Sub0 — declarative JSON backend",
        "PostgreSQL, managed by LingoQL",
        "LingoQL hosting + TLS on 443",
    ]),
    ("SECURITY", [
        "BCRYPT passwords, never returned",
        "Identity from the verified token claim",
        "RBAC enforced in SQL — wrong role, zero rows",
        "Rules writable only while PENDING",
        "Rate limits on auth and mutations",
        "Secrets only via $ENV; audit log INSERT-only",
    ]),
    ("ROADMAP", [
        "Reminders — cron + queues for due dates",
        "SMS / WhatsApp — members live in chat",
        "Payment rails — settle in-app",
        "Portable, signed trust-score export",
        "Receipt uploads as secondary evidence",
        "Multi-currency and diaspora circles",
    ]),
]
cw, gap = Inches(3.83), Inches(0.19)
for i, (name, items) in enumerate(COLS):
    x = M + i * (cw + gap)
    y = Inches(1.86)
    c = card(s, x, y, cw, Inches(3.4), fill="0C3025", line="17513C", soft=False)
    text(s, x + Inches(0.34), y + Inches(0.3), cw - Inches(0.68), Inches(0.3),
         [{"t": name, "size": 9.5, "bold": True, "color": PRIMARY_LIGHT, "font": HEAD, "spc": 1.3}])
    for j, it in enumerate(items):
        iy = y + Inches(0.76) + j * Inches(0.42)
        rect(s, x + Inches(0.36), iy + Inches(0.09), Inches(0.09), Inches(0.09),
             fill=ACCENT if i == 2 else PRIMARY_LIGHT, shape=MSO_SHAPE.OVAL)
        text(s, x + Inches(0.62), iy, cw - Inches(0.96), Inches(0.35),
             [{"t": it, "size": 10, "color": "BFD4C9", "line": 1.3}])

cta = card(s, M, Inches(5.52), CW, Inches(1.16), fill=PRIMARY, line=PRIMARY, soft=False)
shadow(cta, blur=26, dist=10, color="0FA968", alpha=45)
text(s, M + Inches(0.55), Inches(5.76), Inches(8.4), Inches(0.7),
     [{"t": "Keep the tradition. Replace the notebook.", "size": 20, "bold": True,
       "color": WHITE, "font": HEAD, "space_after": 3},
      {"t": "github.com/akaheem/Circle-Safe  ·  README.md · DEPLOYMENT.md · SELF_HOSTING.md · DEMO_SCRIPT.md",
       "size": 9.5, "color": "D8F5E8"}])
pill(s, W - M - Inches(2.5), Inches(5.86), Inches(1.95), Inches(0.48), "CircleSafe",
     fill=WHITE, color=PRIMARY_DARK, size=13)
footer(s, 10, dark=True)

out = "docs/CircleSafe-Documentation.pptx"
prs.save(out)
print(f"saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
      f"fonts: {HEAD} / {SANS})")
